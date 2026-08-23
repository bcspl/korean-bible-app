#!/usr/bin/env python3
"""
Binance-style project status PPTX for Korean Bible App.
Output: Desktop/Bible Project/KoreanBible_ProjectReport_YYYY-MM-DD.pptx
Re-run after each major session to refresh the dated report.
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE

# Binance palette
BG = RGBColor(0x0B, 0x0E, 0x11)
CARD = RGBColor(0x1E, 0x23, 0x29)
CARD2 = RGBColor(0x2B, 0x31, 0x39)
YELLOW = RGBColor(0xF0, 0xB9, 0x0B)
YELLOW_DIM = RGBColor(0xC9, 0x9A, 0x0A)
WHITE = RGBColor(0xEA, 0xEC, 0xF0)
MUTED = RGBColor(0x84, 0x8E, 0x9C)
GREEN = RGBColor(0x0E, 0xCB, 0x81)
RED = RGBColor(0xF6, 0x46, 0x5D)
GRAY_LINE = RGBColor(0x47, 0x4D, 0x57)

W, H = Inches(10), Inches(5.625)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def box(slide, l, t, w, h, text, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Malgun Gothic"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def multi(slide, l, t, w, h, lines, size=11, color=WHITE, bold_first=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Malgun Gothic"
        p.font.bold = bold_first and i == 0
        p.space_after = Pt(4)
    return tb


def card(slide, l, t, w, h, accent=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, CARD if not accent else CARD2)
    sh.adjustments[0] = 0.08
    return sh


def header_bar(slide, title, subtitle=None):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(bg, BG)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
    fill(bar, YELLOW)
    box(slide, Inches(0.45), Inches(0.22), Inches(8), Inches(0.4), title, 22, True, YELLOW)
    if subtitle:
        box(slide, Inches(0.45), Inches(0.55), Inches(9), Inches(0.28), subtitle, 11, False, MUTED)
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.35), W, Inches(0.275))
    fill(foot, CARD)
    box(slide, Inches(0.45), Inches(5.38), Inches(9), Inches(0.22),
        "Korean Bible App  ·  Confidential PM Report  ·  Update after each session", 9, False, MUTED)


def progress_bar(slide, l, t, w, h, pct, label):
    track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(track, CARD2)
    track.adjustments[0] = 0.5
    fw = max(Inches(0.15), w * max(0.02, min(1.0, pct)))
    fill_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, fw, h)
    fill(fill_bar, YELLOW)
    fill_bar.adjustments[0] = 0.5
    box(slide, l, t + h + Inches(0.02), w, Inches(0.22), label, 10, False, MUTED)


def build(out_path: Path):
    today = datetime.now().strftime("%Y-%m-%d")
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ---- 1 Title ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(bg, BG)
    # left yellow strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
    fill(strip, YELLOW)
    box(s, Inches(0.6), Inches(1.3), Inches(8.5), Inches(0.4), "PROJECT STATUS REPORT", 12, True, YELLOW)
    box(s, Inches(0.6), Inches(1.75), Inches(8.8), Inches(0.9), "한국어 성경 앱", 40, True, WHITE)
    box(s, Inches(0.6), Inches(2.7), Inches(8.5), Inches(0.45),
        "비전 · 현재 성과 · 공개 전 과제 · 가치 로드맵", 16, False, MUTED)
    # stats row
    for i, (num, lab) in enumerate([("~92%", "제품 기능"), ("2026-11", "Play Store 목표"), ("3", "PD 역본"), ("102", "PD 찬송")]):
        x = Inches(0.6 + i * 2.25)
        card(s, x, Inches(3.5), Inches(2.05), Inches(1.1))
        box(s, x + Inches(0.12), Inches(3.6), Inches(1.8), Inches(0.5), num, 22, True, YELLOW)
        box(s, x + Inches(0.12), Inches(4.15), Inches(1.8), Inches(0.3), lab, 11, False, MUTED)
    box(s, Inches(0.6), Inches(5.0), Inches(8), Inches(0.3), f"Report date: {today}  ·  Binance design system", 10, False, MUTED)

    # ---- 2 Vision ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "01  비전 & 원칙", "Vision — why this product exists")
    card(s, Inches(0.45), Inches(1.0), Inches(9.1), Inches(1.35))
    multi(s, Inches(0.65), Inches(1.15), Inches(8.7), Inches(1.1), [
        "비전",
        "누구나 저작권 걱정 없이, 완전 오프라인으로 쓰는 한국어 중심 성경·예배 올인원 앱.",
    ], 13, WHITE, True)
    principles = [
        ("PD Only", "개역개정·비-PD 찬송 미수록"),
        ("Offline", "광고·계정 없는 로컬 우선"),
        ("Verified", "본문 검증 스크립트 PASS"),
        ("Practical", "가정·소그룹·개인 묵상"),
    ]
    for i, (t, d) in enumerate(principles):
        x = Inches(0.45 + (i % 4) * 2.35)
        y = Inches(2.6)
        card(s, x, y, Inches(2.2), Inches(1.8))
        box(s, x + Inches(0.15), y + Inches(0.25), Inches(1.9), Inches(0.4), t, 14, True, YELLOW)
        box(s, x + Inches(0.15), y + Inches(0.75), Inches(1.9), Inches(0.8), d, 12, False, MUTED)

    # ---- 3 Status ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "02  현재 상태", "Where we are — product vs store readiness")
    card(s, Inches(0.45), Inches(1.05), Inches(4.4), Inches(3.9))
    box(s, Inches(0.65), Inches(1.2), Inches(4), Inches(0.35), "제품 기능 완성도", 14, True, YELLOW)
    progress_bar(s, Inches(0.65), Inches(1.7), Inches(3.9), Inches(0.28), 0.92, "██████████████████░░  ~92%")
    multi(s, Inches(0.65), Inches(2.25), Inches(3.9), Inches(2.4), [
        "✅ P0–P3 기반·읽기·북마크·탭 100%",
        "✅ 다역본 KRV·KJV·ASV 병렬",
        "✅ 본문 검증 verify PASS",
        "✅ PD 찬송 102 · 교독 51",
        "✅ 주기도문 기도 본문만",
        "✅ flutter test 5/5",
    ], 12, WHITE)

    card(s, Inches(5.1), Inches(1.05), Inches(4.4), Inches(3.9))
    box(s, Inches(5.3), Inches(1.2), Inches(4), Inches(0.35), "Play Store 공개 준비", 14, True, YELLOW)
    progress_bar(s, Inches(5.3), Inches(1.7), Inches(3.9), Inches(0.28), 0.20, "████░░░░░░░░░░░░░░░░  ~20%")
    multi(s, Inches(5.3), Inches(2.25), Inches(3.9), Inches(2.4), [
        "⏳ packageId (com.example 교체)",
        "⏳ 업로드 키스토어 · AAB",
        "⏳ 개인정보처리방침 URL",
        "⏳ 스크린샷·스토어 카피",
        "⏳ 프로덕션 서명 기기 테스트",
        "🔄 APK/Web 재빌드 필요",
    ], 12, WHITE)

    # ---- 4 Done detail ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "03  수행 완료 내역", "Shipped value — already in source")
    rows = [
        ("성경", "KRV 완전본 · KJV · ASV · 절 단위 최대 3역본 병렬 · 검색·북마크"),
        ("품질", "verify_bible_texts PASS · 골든 구절 · 장 누락 수정"),
        ("찬송", "PD 102곡 · 한/영 · 유래·라이선스 배너 · 교육용 SVG"),
        ("예배", "교독 51 · 사도신경/주기도 역본 · 기도 서사 제외"),
        ("UX", "5탭 · 다크/글자크기/고대비 · 스플래시·아이콘"),
        ("문서", "PM 폴더 · 보고서 · 세션 로그 구조"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.0 + i * 0.65)
        card(s, Inches(0.45), y, Inches(9.1), Inches(0.58))
        box(s, Inches(0.6), y + Inches(0.12), Inches(1.3), Inches(0.35), k, 13, True, YELLOW)
        box(s, Inches(2.0), y + Inches(0.12), Inches(7.3), Inches(0.35), v, 12, False, WHITE)

    # ---- 5 Pre-release ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "04  공개 전 필수 과제", "Play Store gate — must finish before 1.0")
    cols = [
        ("정책·고지", ["PD/NKRV 미수록 문구", "개인정보 URL", "연령 등급", "지원 이메일"]),
        ("Android 기술", ["applicationId 교체", "키스토어", "AAB 빌드", "targetSdk 충족"]),
        ("품질", ["비행기모드 스모크", "저사양 로딩", "서명 APK 설치", "릴리스 재빌드"]),
        ("스토어 자산", ["스크린샷 4–8", "피처 그래픽", "짧은/긴 설명", "아이콘 최종"]),
    ]
    for i, (title, items) in enumerate(cols):
        x = Inches(0.4 + i * 2.4)
        card(s, x, Inches(1.05), Inches(2.25), Inches(3.9))
        box(s, x + Inches(0.12), Inches(1.2), Inches(2.0), Inches(0.35), title, 13, True, YELLOW)
        multi(s, x + Inches(0.12), Inches(1.7), Inches(2.0), Inches(3.0),
              [f"• {it}" for it in items], 11, WHITE)

    # ---- 6 Timeline ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "05  일정 — 2026-11 Play Store", "Milestone roadmap to public launch")
    months = [
        ("8월", "콘텐츠 동결 후보\nPM 구조\n재빌드", GREEN),
        ("9월", "packageId\n키스토어·AAB\n정책 문구", YELLOW),
        ("10월", "기기 테스트\n온보딩·복사\n스크린샷", YELLOW),
        ("11월", "스토어 제출\n심사 대응\n1.0 오픈", YELLOW),
    ]
    # timeline line
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.55), Inches(8.4), Inches(0.06))
    fill(line, GRAY_LINE)
    for i, (m, body, col) in enumerate(months):
        x = Inches(0.7 + i * 2.3)
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), Inches(2.4), Inches(0.35), Inches(0.35))
        fill(oval, col)
        card(s, x, Inches(1.15), Inches(2.1), Inches(1.05))
        box(s, x + Inches(0.1), Inches(1.25), Inches(1.9), Inches(0.35), m, 16, True, YELLOW, PP_ALIGN.CENTER)
        card(s, x, Inches(3.0), Inches(2.1), Inches(1.85))
        multi(s, x + Inches(0.12), Inches(3.15), Inches(1.85), Inches(1.6), body.split("\n"), 12, WHITE)
    box(s, Inches(0.5), Inches(5.0), Inches(9), Inches(0.25),
        "M4 스토어 기술(9) → M5 공개 품질(10) → M6 Play 오픈(11)", 11, False, MUTED)

    # ---- 7 Value ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "06  사용성 · 가치 로드맵", "What users gain — prioritized backlog")
    # two columns
    card(s, Inches(0.45), Inches(1.05), Inches(4.4), Inches(3.9))
    box(s, Inches(0.65), Inches(1.2), Inches(4), Inches(0.35), "P0  공개 전 (가치 높음)", 13, True, YELLOW)
    multi(s, Inches(0.65), Inches(1.7), Inches(4.0), Inches(3.0), [
        "V1 온보딩 — 첫인상·PD 이해",
        "V2 역본 선택 저장 — 마찰 제거",
        "V3 절 복사/공유 — 소그룹 확산",
        "V4 라이선스 한눈에 — 신뢰·CS",
        "",
        "→ RELEASE 체크리스트 R22–R25",
    ], 12, WHITE)

    card(s, Inches(5.1), Inches(1.05), Inches(4.4), Inches(3.9))
    box(s, Inches(5.3), Inches(1.2), Inches(4), Inches(0.35), "P1–P2  공개 후 / 선택", 13, True, YELLOW)
    multi(s, Inches(5.3), Inches(1.7), Inches(4.0), Inches(3.0), [
        "위치 복원 · 통독 계획 · 하이라이트",
        "교독 인도자/회중 강조 모드",
        "태블릿 2열 · TTS · 찬송 필터",
        "",
        "비권장: 개역개정 · 비-PD 찬송",
        "          강제 계정/클라우드",
    ], 12, WHITE)

    # ---- 8 PM structure ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "07  프로젝트 관리 구조", "Update after every work session")
    card(s, Inches(0.45), Inches(1.05), Inches(9.1), Inches(1.5))
    multi(s, Inches(0.65), Inches(1.2), Inches(8.7), Inches(1.2), [
        "docs/project-management/",
        "README · VISION_AND_PLAN · RELEASE_CHECKLIST_PLAYSTORE",
        "VALUE_AND_USABILITY · SESSION_UPDATE_LOG",
    ], 13, WHITE, True)

    steps = [
        ("1", "SESSION_LOG", "오늘 한 일 / 막힌 일 / 다음 1건"),
        ("2", "VALUE", "완료 항목 ✅ 체크"),
        ("3", "RELEASE", "공개 전 상태 변경"),
        ("4", "VISION", "진행률 % 한 줄"),
        ("5", "PPTX", "날짜 파일명 재생성 (선택)"),
    ]
    for i, (n, t, d) in enumerate(steps):
        x = Inches(0.45 + i * 1.9)
        card(s, x, Inches(2.8), Inches(1.8), Inches(2.1))
        box(s, x + Inches(0.1), Inches(2.95), Inches(1.6), Inches(0.35), n, 18, True, YELLOW, PP_ALIGN.CENTER)
        box(s, x + Inches(0.1), Inches(3.4), Inches(1.6), Inches(0.35), t, 11, True, WHITE, PP_ALIGN.CENTER)
        box(s, x + Inches(0.1), Inches(3.85), Inches(1.6), Inches(0.8), d, 10, False, MUTED, PP_ALIGN.CENTER)

    # ---- 9 Risks ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "08  리스크 & 성공 지표", "Stay green until November")
    risks = [
        ("packageId 미변경", "Play 업로드 거부", "9월 1주 내 확정"),
        ("키스토어 분실", "업데이트 불가", "오프라인 금고 백업"),
        ("정책 문구 미비", "심사 반려", "PD/수집없음 명시"),
        ("대용량 첫 로딩", "이탈", "스플래시·진행 표시"),
    ]
    for i, (a, b, c) in enumerate(risks):
        y = Inches(1.05 + i * 0.85)
        card(s, Inches(0.45), y, Inches(9.1), Inches(0.75))
        box(s, Inches(0.65), y + Inches(0.2), Inches(2.5), Inches(0.4), a, 13, True, RED)
        box(s, Inches(3.3), y + Inches(0.2), Inches(3.0), Inches(0.4), b, 12, False, WHITE)
        box(s, Inches(6.5), y + Inches(0.2), Inches(2.8), Inches(0.4), c, 12, False, YELLOW)

    # ---- 10 Next ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "09  다음 액션 · Top 3", "Start here next session")
    tops = [
        ("01", "스토어 기술 게이트", "packageId · 키스토어 · AAB\nprivacy URL · 지원 이메일"),
        ("02", "공개 최소 UX", "온보딩 · 역본 저장\n절 복사 · 라이선스 화면"),
        ("03", "품질·제출", "재빌드 · 기기 스모크\n스크린샷 · Play 제출"),
    ]
    for i, (n, t, d) in enumerate(tops):
        x = Inches(0.5 + i * 3.15)
        card(s, x, Inches(1.2), Inches(3.0), Inches(3.5), accent=True)
        box(s, x + Inches(0.2), Inches(1.45), Inches(2.6), Inches(0.45), n, 28, True, YELLOW)
        box(s, x + Inches(0.2), Inches(2.1), Inches(2.6), Inches(0.5), t, 16, True, WHITE)
        multi(s, x + Inches(0.2), Inches(2.8), Inches(2.6), Inches(1.5), d.split("\n"), 13, MUTED)
    box(s, Inches(0.5), Inches(4.9), Inches(9), Inches(0.3),
        f"목표: 2026-11 Play Store  ·  보고서 파일: KoreanBible_ProjectReport_{today}.pptx", 11, False, YELLOW)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main():
    today = datetime.now().strftime("%Y-%m-%d")
    desktop = Path.home() / "Desktop" / "Bible Project"
    name = f"KoreanBible_ProjectReport_{today}.pptx"
    out = desktop / name
    # also copy to project docs archive
    proj = Path(__file__).resolve().parents[2]
    archive = proj / "docs" / "roadmap-pptx" / name
    p = build(out)
    archive.parent.mkdir(parents=True, exist_ok=True)
    import shutil
    shutil.copy2(p, archive)
    print(f"OK Desktop: {p}")
    print(f"OK Archive: {archive}")


if __name__ == "__main__":
    main()
