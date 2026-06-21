#!/usr/bin/env python3
"""
Korean Bible App - Roadmap and Task Proposal Presentation
Integrates the full roadmap and prioritized task backlog.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Color palette - Calm Navy + Teal for Bible/Tech theme
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
TEAL = RGBColor(0x0D, 0x94, 0x88)
GOLD = RGBColor(0xD9, 0x77, 0x06)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)

def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    set_shape_fill(background, NAVY)
    background.line.fill.background()
    
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(0.15), Inches(1.2))
    set_shape_fill(bar, TEAL)
    bar.line.fill.background()
    
    # Title
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1),
                 "한국어 성경 앱", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    # Subtitle
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(0.8),
                 "로드맵 및 작업 제안", font_size=28, bold=False, color=RGBColor(0xCC, 0xFB, 0xF1), align=PP_ALIGN.CENTER)
    
    # Date and version
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(0.5),
                 "Updated 2026-06-20 | 완전 무료 오프라인 한국어 성경 앱", font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
    
    return slide

def add_content_slide(prs, title, content_lines, accent_color=TEAL):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
    set_shape_fill(top_bar, accent_color)
    top_bar.line.fill.background()
    
    # Title
    title_box = add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7), title, font_size=28, bold=True, color=NAVY)
    
    # Content
    y_pos = 1.1
    for line in content_lines:
        if line.startswith("##"):
            add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4), line[2:].strip(), font_size=18, bold=True, color=TEAL)
            y_pos += 0.45
        elif line.startswith("-"):
            add_text_box(slide, Inches(0.7), Inches(y_pos), Inches(8.8), Inches(0.35), "• " + line[1:].strip(), font_size=13, color=DARK_TEXT)
            y_pos += 0.32
        elif line.strip():
            add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.35), line, font_size=13, color=DARK_TEXT)
            y_pos += 0.35
        else:
            y_pos += 0.15
    
    return slide

def add_phases_overview(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
    set_shape_fill(top_bar, TEAL)
    top_bar.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.6), "개발 단계 개요 (Phases)", font_size=26, bold=True, color=NAVY)
    
    phases = [
        ("Phase 0", "기반 구축", "완료", NAVY),
        ("Phase 1", "데이터 계층", "최우선", TEAL),
        ("Phase 2", "핵심 읽기 MVP", "MVP", TEAL),
        ("Phase 3", "사용자 기능", "중요", GOLD),
        ("Phase 4", "UI/UX 고도화", "품질", GOLD),
        ("Phase 5", "배포 & 플랫폼", "릴리스", NAVY),
    ]
    
    x_start = 0.4
    card_w = 1.45
    card_h = 2.8
    y = 1.1
    
    for i, (phase, desc, status, color) in enumerate(phases):
        x = x_start + i * (card_w + 0.08)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        
        # Phase number header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(0.5))
        set_shape_fill(header, color)
        header.line.fill.background()
        
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.08), Inches(card_w - 0.1), Inches(0.4), phase, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.6), Inches(card_w - 0.1), Inches(0.6), desc, font_size=10, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        
        # Status
        status_color = TEAL if status == "완료" or "최우선" in status or "MVP" in status else GOLD
        add_text_box(slide, Inches(x + 0.05), Inches(y + 1.4), Inches(card_w - 0.1), Inches(0.35), status, font_size=9, bold=True, color=status_color, align=PP_ALIGN.CENTER)
    
    # Description below
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(1.2), 
                 "Phase 1부터 시작하여 데이터 계층을 완성한 후, 순차적으로 MVP → 기능 → 품질 → 배포로 진행합니다. 각 Phase의 상세 작업은 다음 슬라이드 참조.",
                 font_size=12, color=DARK_TEXT)
    
    return slide

def add_phase_detail(prs, title, items, phase_num):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
    set_shape_fill(top_bar, TEAL if phase_num <= 2 else GOLD)
    top_bar.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.6), title, font_size=24, bold=True, color=NAVY)
    
    y = 1.0
    for item in items:
        # Bullet card style
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y + 0.05), Inches(0.15), Inches(0.15))
        set_shape_fill(bullet, TEAL)
        bullet.line.fill.background()
        
        add_text_box(slide, Inches(0.8), Inches(y), Inches(8.7), Inches(0.45), item, font_size=13, color=DARK_TEXT)
        y += 0.48
    
    return slide

def add_task_table(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
    set_shape_fill(top_bar, GOLD)
    top_bar.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.5), "우선순위 작업 순서 (요약)", font_size=24, bold=True, color=NAVY)
    
    # Simple text based prioritized list
    tasks = [
        "1. Phase 1: 데이터 모델 정의 (BibleBook, Chapter, Verse)",
        "2. Hive 어댑터 생성 및 build_runner 실행",
        "3. JSON 파서 + 앱 시작 시 데이터 로드 구현",
        "4. Phase 2: 책 목록 + 장 선택 + 본문 뷰어 UI",
        "5. 네비게이션 및 현재 위치 저장",
        "6. 북마크 / 검색 기능 (Hive 기반)",
        "7. 설정 화면 및 UI 고도화",
        "8. Android/iOS 빌드 검증 및 릴리스 준비"
    ]
    
    y = 0.95
    for task in tasks:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.42))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        
        add_text_box(slide, Inches(0.7), Inches(y + 0.05), Inches(8.6), Inches(0.35), task, font_size=12, color=DARK_TEXT)
        y += 0.5
    
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.4), 
                 "상세 백로그: docs/Task_Prioritized_Backlog_and_Schedule_2026-06-19.md 참조",
                 font_size=10, color=RGBColor(0x64, 0x74, 0x8B))
    
    return slide

def add_risks_and_next(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
    set_shape_fill(top_bar, NAVY)
    top_bar.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.5), "리스크 및 다음 단계", font_size=24, bold=True, color=NAVY)
    
    # Left column: Risks
    add_text_box(slide, Inches(0.5), Inches(0.9), Inches(4.3), Inches(0.4), "주요 리스크", font_size=16, bold=True, color=TEAL)
    risks = [
        "• 대용량 JSON 데이터 (7MB+) 로딩 성능",
        "• 모바일에서의 메모리/성능 최적화",
        "• 데이터 정확성 검증 필요",
        "• 문서(백로그/스펙) 지속 업데이트"
    ]
    y = 1.35
    for r in risks:
        add_text_box(slide, Inches(0.5), Inches(y), Inches(4.3), Inches(0.4), r, font_size=12, color=DARK_TEXT)
        y += 0.38
    
    # Right column: Next
    add_text_box(slide, Inches(5.2), Inches(0.9), Inches(4.3), Inches(0.4), "즉시 추천 작업", font_size=16, bold=True, color=GOLD)
    nexts = [
        "1. 데이터 모델 + Hive 어댑터 (Phase 1)",
        "2. JSON 로더 구현",
        "3. 기본 성경 읽기 UI (Phase 2)",
        "4. Windows 앱 실행 테스트"
    ]
    y = 1.35
    for n in nexts:
        add_text_box(slide, Inches(5.2), Inches(y), Inches(4.3), Inches(0.4), n, font_size=12, color=DARK_TEXT)
        y += 0.38
    
    # Bottom call to action
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(9), Inches(1.4))
    set_shape_fill(cta_box, RGBColor(0xCC, 0xFB, 0xF1))
    cta_box.line.fill.background()
    
    add_text_box(slide, Inches(0.7), Inches(4.0), Inches(8.6), Inches(1.0),
                 "로드맵 문서: docs/Project_Status_and_Updated_Roadmap_2026-06-19.md\n" +
                 "백로그 문서: docs/Task_Prioritized_Backlog_and_Schedule_2026-06-19.md\n\n" +
                 "이 문서를 기반으로 Phase 1부터 순차 진행을 권장합니다.",
                 font_size=12, color=NAVY)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Overview
    add_content_slide(prs, "프로젝트 개요", [
        "## 비전",
        "완전 무료 오프라인 한국어 성경 앱 (Public Domain - 개역개정 등)",
        "",
        "## 주요 목표",
        "- 인터넷 없이 성경 읽기, 검색, 북마크, 메모 지원",
        "- 접근성 중심 UI (큰 글씨, 다크모드)",
        "- Flutter로 모든 플랫폼 지원 (Android, iOS, Web, Desktop)",
        "",
        "## 데이터",
        "Public Domain 한국어 성경 데이터 활용"
    ])
    
    # Slide 3: Current Status
    add_content_slide(prs, "현재 상태 (2026-06-20)", [
        "## 완료",
        "- Flutter 프로젝트 초기화 및 다중 플랫폼 준비",
        "- Hive, Provider 등 의존성 설정 완료",
        "- 성경 데이터 assets/data/에 배치",
        "- 기본 앱 스켈레톤 (Hive init + 한국어 UI)",
        "- flutter analyze 클린 + Windows 빌드 성공",
        "",
        "## 진행 중 / 미완성",
        "- 실제 성경 데이터 로딩 및 UI 미구현",
        "- Hive 모델/어댑터 미구현",
        "- 로드맵/백로그 문서 이제 작성 완료"
    ])
    
    # Slide 4: Tech Stack
    add_content_slide(prs, "기술 스택", [
        "## 핵심",
        "- Flutter (크로스 플랫폼)",
        "- Hive (오프라인 NoSQL DB) + hive_generator",
        "- Provider (상태 관리)",
        "",
        "## UI & 기타",
        "- Material Design + 접근성 (폰트 크기)",
        "- Assets: assets/data/ (성경 JSON)",
        "",
        "## 향후 고려",
        "- shared_preferences (설정 저장)",
        "- 데이터 최적화 (lazy load, 인덱싱)"
    ])
    
    # Slide 5: Phases Overview
    add_phases_overview(prs)
    
    # Slide 6: Phase 1 Detail
    add_phase_detail(prs, "Phase 1: 데이터 계층 (최우선)", [
        "성경 데이터 모델 정의 (BibleBook, Chapter, Verse, VerseRef)",
        "Hive 어댑터 생성 (@HiveType, @HiveField)",
        "build_runner 실행으로 어댑터 코드 생성",
        "JSON 파서 서비스 구현 (rootBundle 사용)",
        "앱 시작 시 데이터 로드 + Hive Box 시드",
        "데이터 검증 및 로딩 성능 측정"
    ], 1)
    
    # Slide 7: Phase 2 Detail
    add_phase_detail(prs, "Phase 2: 핵심 읽기 MVP", [
        "책 목록 화면 (구약/신약 분류)",
        "장 선택 화면 (그리드 또는 목록)",
        "본문 뷰어 (절 단위 표시, 스크롤)",
        "이전/다음 장·책 네비게이션",
        "현재 읽기 위치 저장 (Provider + 로컬)",
        "하단 네비게이션 연결 (성경/북마크/검색)"
    ], 2)
    
    # Slide 8: Phase 3+ and Tasks
    add_phase_detail(prs, "Phase 3+: 사용자 기능 및 고도화", [
        "북마크, 하이라이트, 메모 (Hive 저장)",
        "본문 검색 기능",
        "설정 화면 (글꼴 크기, 테마, 버전)",
        "다크모드 완성 + 접근성 강화",
        "Android / iOS / Web 빌드 검증",
        "앱 아이콘, 릴리스 준비 및 스토어 제출"
    ], 3)
    
    # Slide 9: Prioritized Tasks
    add_task_table(prs)
    
    # Slide 10: Risks + Next Steps
    add_risks_and_next(prs)
    
    # Save
    output_path = "docs/Korean_Bible_App_Roadmap_및_작업제안.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()